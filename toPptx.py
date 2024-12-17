from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation object
presentation = Presentation()

# Helper function to add a slide
def add_slide(title, content, image_path=None):
    slide_layout = presentation.slide_layouts[1]
    slide = presentation.slides.add_slide(slide_layout)
    
    title_placeholder = slide.shapes.title
    content_placeholder = slide.placeholders[1]
    
    title_placeholder.text = title
    content_placeholder.text = content
    
    if image_path:
        slide.shapes.add_picture(image_path, Inches(5.5), Inches(1.5), height=Inches(3))
        
    # Formatting title text
    title_placeholder.text_frame.paragraphs[0].font.size = Pt(30)
    title_placeholder.text_frame.paragraphs[0].font.bold = True
    title_placeholder.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    title_placeholder.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

# Slide 1: Title Slide
add_slide("Health Monitoring and Inhaler Tracking System", 
          "A Dual Device System for Continuous Health Monitoring and Medical Equipment Accessibility\n\nPresenters: Team\nDate: Institution")

# Slide 2: Abstract
add_slide("Abstract",
          "This project presents a wearable health monitoring system that records real-time health metrics "
          "and transmits data to a Firebase cloud for analysis by healthcare providers. It includes a beacon attachable "
          "to inhalers, enabling asthmatic patients to prevent loss of their inhalers. Our design enhances healthcare "
          "accessibility and patient safety.")

# Slide 3: Problem Statement
add_slide("Problem Statement",
          "Many patients face challenges in maintaining continuous health monitoring and tracking essential medical items "
          "like inhalers, leading to emergencies. Existing devices lack integrated systems for both real-time health tracking "
          "and item location awareness. This project addresses these issues, offering a combined health monitoring and tracking solution.")

# Slide 4: Solution Overview
add_slide("Solution Overview",
          "Our project combines a wearable device for health monitoring with a beacon attachable to an inhaler, "
          "providing continuous tracking and secure, real-time data access for healthcare providers. This dual approach enhances "
          "patient safety by keeping health data accessible and inhalers within reach.")

# Slide 5: System Architecture
add_slide("System Architecture",
          "The wearable device includes an ESP32 microcontroller, MAX30102 pulse oximeter, and DS18B20 temperature sensor, "
          "while the beacon connects to the inhaler to prevent misplacement. This architecture ensures seamless data flow "
          "and enhances health monitoring and device tracking.")

# Slide 6: Flowchart
add_slide("Flowchart",
          "This flowchart illustrates each step: data collection, transmission to Firebase, storage, analysis, alert system, "
          "and user/doctor interface. Each part contributes to an efficient monitoring and tracking system, increasing accessibility and response speed.")

# Slide 7: Device Design and Components
add_slide("Device Design and Components",
          "ESP32 handles processing and data transmission; MAX30102 monitors heart rate and oxygen level; DS18B20 measures body temperature; "
          "and the beacon attaches to inhalers, sending alerts if the inhaler is misplaced. Together, they create a compact, "
          "reliable health monitoring system.")

# Slide 8: Methodology
add_slide("Methodology",
          "The methodology covers device construction, data transmission, user interface design, and testing. Data from ESP32 is transmitted to Firebase, "
          "alerts are generated for unusual readings, and users are alerted if their inhaler is out of range.")

# Slide 9: Societal Impact & Economic Advantage
add_slide("Societal Impact & Economic Advantage",
          "This project enhances patient safety and reduces health risks associated with lost inhalers while providing an economic advantage "
          "by minimizing emergency healthcare costs and hospital visits. It supports proactive healthcare management with potential insurance incentives.")

# Slide 10: Conclusion
add_slide("Conclusion",
          "This project meets the need for remote health monitoring and secure, accessible devices for emergencies. It improves healthcare access, "
          "particularly for individuals with chronic illnesses, offering timely diagnosis and reducing emergency risks.")

# Save the presentation
pptx_file_path = "Health_Monitoring_Project_Presentation.pptx"
presentation.save(pptx_file_path)

pptx_file_path
